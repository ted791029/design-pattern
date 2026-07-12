#!/usr/bin/env python3
"""
Puzzle Tree Presentation Generator
Creates 3 slides showing progressive assembly of a puzzle-piece tree:
  Slide 1: 萌芽 — only the root/base puzzle piece
  Slide 2: 成長 — roots + trunk pieces
  Slide 3: 茁壯 — complete tree (all three pieces)

Slide layout mirrors the THUS slides from the original presentation.
"""

import io
import os
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ─────────────────────────────────────────────
# Font setup (CJK)
# ─────────────────────────────────────────────
_FONT_CANDIDATES = [
    '/tmp/NotoSansTC-Regular.otf',
    '/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf',
]
CJK_PROP = None
for fp in _FONT_CANDIDATES:
    if os.path.exists(fp):
        CJK_PROP = fm.FontProperties(fname=fp)
        break

def text_kw(size=14, color='white', bold=False):
    kw = dict(fontsize=size, color=color, fontweight='bold' if bold else 'normal')
    if CJK_PROP:
        kw['fontproperties'] = CJK_PROP
    return kw

# ─────────────────────────────────────────────
# Color constants
# ─────────────────────────────────────────────

# PowerPoint RGB
NAVY       = RGBColor(0x1A, 0x29, 0x6B)
DARK_NAVY  = RGBColor(0x0F, 0x1E, 0x52)
GOLD_RGB   = RGBColor(0xC9, 0xA0, 0x2B)
WHITE_RGB  = RGBColor(0xFF, 0xFF, 0xFF)
MAROON_RGB = RGBColor(0x6B, 0x1A, 0x2B)
GREY_RGB   = RGBColor(0x55, 0x55, 0x77)

# Stage accent colors (PowerPoint)
ACCENT = {
    '萌芽': RGBColor(0xC1, 0x7F, 0x3E),   # warm amber-brown
    '成長': RGBColor(0x3D, 0x8A, 0x3A),   # medium green
    '茁壯': RGBColor(0x1B, 0x6E, 0x3E),   # deep forest green
}

# Matplotlib hex colors for tree pieces
ROOT_C   = '#8B5E3C'   # warm brown
TRUNK_C  = '#3D7A2A'   # medium green
CANOPY_C = '#1E5E1E'   # dark forest green

# ─────────────────────────────────────────────
# Puzzle geometry helpers
# ─────────────────────────────────────────────

def arc(cx, cy, r, t0, t1, n=35):
    """Return (xs, ys) arrays for an arc from angle t0 to t1 (radians)."""
    ts = np.linspace(t0, t1, n)
    return cx + r * np.cos(ts), cy + r * np.sin(ts)


def tab_up(cx, base_y, r, n=35):
    """
    Semicircular tab pointing UP at (cx, base_y).
    Traversal: right-to-left (theta 0 → π), arc goes above base_y.
    Returns xs going from cx+r to cx-r while rising.
    """
    return arc(cx, base_y, r, 0, np.pi, n)


def blank_down(cx, base_y, r, n=35):
    """
    Semicircular blank (indentation) pointing DOWN at (cx, base_y).
    Traversal: left-to-right (theta π → 2π), arc dips below base_y.
    Returns xs going from cx-r to cx+r while dipping.
    """
    return arc(cx, base_y, r, np.pi, 2 * np.pi, n)


# ─────────────────────────────────────────────
# Puzzle piece path builders
# ─────────────────────────────────────────────
# Coordinate system: xlim 0–10, ylim 0–14

TAB_R = 0.40   # radius of puzzle tab / blank semi-circle
TAB_CX = 5.0   # x-center of all connectors (vertical tree, center x=5)


def roots_path():
    """
    Trapezoid for the roots/base piece.
    Bottom: y=0.3, x from 1.0 to 9.0  (wide)
    Top:    y=3.3, x from 3.5 to 6.5  (narrow, matches trunk width)
    Tab UP on the top edge (connects to trunk blank).
    Traversal: clockwise → bottom-L, bottom-R, top-R, [tab right→left going up], top-L, back.
    """
    r = TAB_R
    cx = TAB_CX
    y0, y1 = 0.3, 3.3
    xl0, xr0 = 1.0, 9.0   # wide bottom
    xl1, xr1 = 3.5, 6.5   # narrow top (trunk width)

    tx, ty = tab_up(cx, y1, r)   # right→left going UP

    xs = [xl0, xr0, xr1, cx + r] + list(tx) + [xl1, xl0]
    ys = [y0,  y0,  y1,  y1    ] + list(ty) + [y1,  y0 ]
    return xs, ys


def trunk_path():
    """
    Rectangle for the trunk piece.
    Bottom: y=3.3, x from 3.5 to 6.5
    Top:    y=6.3, x from 3.5 to 6.5
    Blank DOWN on bottom (receives roots tab).
    Tab  UP  on top    (connects to canopy blank).
    """
    r = TAB_R
    cx = TAB_CX
    y0, y1 = 3.3, 6.3
    xl, xr = 3.5, 6.5

    bx, by = blank_down(cx, y0, r)   # left→right dipping DOWN
    tx, ty = tab_up   (cx, y1, r)   # right→left going UP

    xs = ([xl, cx - r] + list(bx) + [xr, xr, cx + r]
          + list(tx) + [xl, xl])
    ys = ([y0, y0    ] + list(by) + [y0, y1, y1    ]
          + list(ty) + [y1, y0])
    return xs, ys


def canopy_path():
    """
    Circular crown for the canopy piece.
    Base: y=6.3 (connects to trunk tab)
    Circle: center (5, 9.5), radius 4.0  →  top at y≈13.5
    Blank DOWN on base edge (receives trunk tab).
    Upper boundary: major arc of the circle.
    """
    r = TAB_R
    cx = TAB_CX
    y_base = 6.3

    ccx, ccy, cr = 5.0, 9.5, 4.0  # canopy circle

    # Find angles where the circle intersects y=y_base
    dy = y_base - ccy                      # = 6.3 - 9.5 = -3.2
    sin_val = np.clip(dy / cr, -1, 1)      # = -0.8
    t_right = np.arcsin(sin_val)           # ≈ -53.1° ≈ -0.927 rad  (lower-right)
    t_left  = np.pi - t_right              # ≈ π+0.927 ≈ 4.069 rad  (lower-left)

    cx_right = ccx + cr * np.cos(t_right)  # ≈ 5 + 4*0.6 = 7.4
    cx_left  = ccx + cr * np.cos(t_left)   # ≈ 5 - 4*0.6 = 2.6

    # Blank at base (left → right, dipping down)
    bx, by = blank_down(cx, y_base, r)

    # Upper arc: from t_right going counter-clockwise (increasing θ) to t_left
    # This sweeps OVER THE TOP of the circle.
    ax_arr, ay_arr = arc(ccx, ccy, cr, t_right, t_left, n=90)

    xs = ([cx_left, cx - r] + list(bx) + [cx_right]
          + list(ax_arr) + [cx_left])
    ys = ([y_base,  y_base ] + list(by) + [y_base  ]
          + list(ay_arr) + [y_base ])
    return xs, ys


# ─────────────────────────────────────────────
# Tree image renderer
# ─────────────────────────────────────────────

def draw_tree(show_roots=True, show_trunk=False, show_canopy=False):
    """
    Render the puzzle tree and return a BytesIO PNG (transparent background).
    figsize=(4.0, 5.6) with xlim=(0,10), ylim=(0,14) → 1 unit = 0.4 inches.
    """
    FIG_W, FIG_H = 4.0, 5.6   # inches  (ratio = 10/14 = 0.714)
    fig = plt.figure(figsize=(FIG_W, FIG_H))
    ax  = fig.add_axes([0, 0, 1, 1])   # full-figure axes (no margins)
    fig.patch.set_alpha(0)
    ax.set_facecolor('none')
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 14)
    ax.axis('off')

    lw_border = 3.5   # white border linewidth

    def draw_piece(xs, ys, color):
        ax.fill(xs, ys, color=color, zorder=2)
        xs_c = xs + [xs[0]]
        ys_c = ys + [ys[0]]
        ax.plot(xs_c, ys_c, '-', color='white', lw=lw_border,
                zorder=3, solid_capstyle='round', solid_joinstyle='round')

    if show_roots:
        xs, ys = roots_path()
        draw_piece(xs, ys, ROOT_C)
        ax.text(5.0, 1.8, '萌芽', ha='center', va='center',
                zorder=4, **text_kw(size=18, color='white', bold=True))

    if show_trunk:
        xs, ys = trunk_path()
        draw_piece(xs, ys, TRUNK_C)
        ax.text(5.0, 4.8, '成長', ha='center', va='center',
                zorder=4, **text_kw(size=18, color='white', bold=True))

    if show_canopy:
        xs, ys = canopy_path()
        draw_piece(xs, ys, CANOPY_C)
        ax.text(5.0, 9.5, '茁壯', ha='center', va='center',
                zorder=4, **text_kw(size=22, color='white', bold=True))

    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=150, transparent=True,
                bbox_inches=None)
    buf.seek(0)
    plt.close(fig)
    return buf


# ─────────────────────────────────────────────
# Slide content data
# ─────────────────────────────────────────────

SLIDE_DATA = [
    {
        'stage':     '萌芽',
        'sub':       '建立連結',
        'bullets': [
            '與國際夥伴建立合作關係',
            '開啟學術交流與人才互動',
            '奠定信任與合作基礎',
        ],
        'bottom':    '每一個連結，都是改變的開始',
        'show_roots':  True,
        'show_trunk':  False,
        'show_canopy': False,
        'page': 13,
    },
    {
        'stage':     '成長',
        'sub':       '深化合作',
        'bullets': [
            '共同研究與計畫合作',
            '資源共享與能力互補',
            '拓展國際能見度與影響力',
        ],
        'bottom':    '深化合作，讓影響力持續擴大',
        'show_roots':  True,
        'show_trunk':  True,
        'show_canopy': False,
        'page': 14,
    },
    {
        'stage':     '茁壯',
        'sub':       '共創價值',
        'bullets': [
            '解決全球議題',
            '培育國際領袖',
            '提升全球影響力',
            '實現永續發展',
        ],
        'bottom':    '共創價值，成就永續影響力',
        'show_roots':  True,
        'show_trunk':  True,
        'show_canopy': True,
        'page': 15,
    },
]

QUOTE = ('"Well-managed partnership produces measurable academic '
         'excellence, which translates into university international branding"')

TITLE = '國際網絡經營至品牌國際影響力'

# ─────────────────────────────────────────────
# PowerPoint helpers
# ─────────────────────────────────────────────

def inches(v):
    return Inches(v)


def add_rect(slide, x, y, w, h, fill_rgb=None, line=False):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        inches(x), inches(y), inches(w), inches(h)
    )
    shape.line.fill.background()
    if not line:
        shape.line.fill.background()
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_size=12, bold=False,
                italic=False, color=WHITE_RGB, align=PP_ALIGN.LEFT,
                font_name='微軟正黑體', wrap=True, space_before=0):
    txBox = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_para(tf, text, font_size=12, bold=False, italic=False,
             color=WHITE_RGB, align=PP_ALIGN.LEFT,
             font_name='微軟正黑體', space_before=4):
    """Append a paragraph to an existing TextFrame."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return p


# ─────────────────────────────────────────────
# Slide builder
# ─────────────────────────────────────────────

def build_slide(prs, data):
    slide_layout = prs.slide_layouts[6]   # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # ── Dimensions ──────────────────────────────
    SW = 13.33   # slide width  (inches)
    SH = 7.5     # slide height (inches)
    LP = 3.3     # left panel width
    FT = 0.38    # footer bar height
    FY = SH - FT # footer bar top y

    # ── Left panel background ────────────────────
    add_rect(slide, 0, 0, LP, SH, fill_rgb=NAVY)

    # ── Left panel: "THUS" ───────────────────────
    thus_box = slide.shapes.add_textbox(
        inches(0.18), inches(0.28), inches(LP - 0.25), inches(1.1))
    tf = thus_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'THUS'
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.italic = True
    run.font.color.rgb = GOLD_RGB
    run.font.name = 'Impact'

    # Gold divider line under THUS
    line_shape = slide.shapes.add_shape(
        1, inches(0.2), inches(1.42), inches(LP - 0.4), inches(0.04))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = GOLD_RGB
    line_shape.line.fill.background()

    # Left panel subtitle
    add_textbox(
        slide, 0.18, 1.52, LP - 0.28, 1.3,
        '從夥伴經營到品牌成長',
        font_size=16, bold=True, color=WHITE_RGB,
        align=PP_ALIGN.LEFT, wrap=True
    )

    # Left panel bottom: lightbulb + motto
    bulb_box = slide.shapes.add_textbox(
        inches(0.18), inches(5.3), inches(LP - 0.22), inches(1.25))
    tf_b = bulb_box.text_frame
    tf_b.word_wrap = True
    p_b = tf_b.paragraphs[0]
    run_b = p_b.add_run()
    run_b.text = '💡'
    run_b.font.size = Pt(20)
    run_b.font.name = 'Segoe UI Emoji'

    add_para(tf_b, '共創學術價值', font_size=14, bold=True,
             color=WHITE_RGB, align=PP_ALIGN.LEFT, space_before=2)
    add_para(tf_b, '放大國際影響', font_size=14, bold=True,
             color=WHITE_RGB, align=PP_ALIGN.LEFT, space_before=1)

    # ── Footer bar ───────────────────────────────
    add_rect(slide, 0, FY, SW, FT, fill_rgb=MAROON_RGB)
    add_textbox(
        slide, 0.12, FY + 0.05, 6.5, FT - 0.08,
        '臺北醫學大學  ·  Office of Global Engagement',
        font_size=9, color=WHITE_RGB, align=PP_ALIGN.LEFT,
        font_name='微軟正黑體'
    )
    add_textbox(
        slide, SW - 0.6, FY + 0.04, 0.55, FT - 0.06,
        str(data['page']),
        font_size=11, bold=True, color=WHITE_RGB,
        align=PP_ALIGN.RIGHT, font_name='Calibri'
    )

    # ── Right panel ──────────────────────────────
    # Title
    title_box = slide.shapes.add_textbox(
        inches(LP + 0.15), inches(0.1),
        inches(SW - LP - 0.3), inches(0.65)
    )
    tf_t = title_box.text_frame
    p_t = tf_t.paragraphs[0]
    p_t.alignment = PP_ALIGN.LEFT
    run_t = p_t.add_run()
    run_t.text = TITLE
    run_t.font.size = Pt(18)
    run_t.font.bold = True
    run_t.font.color.rgb = NAVY
    run_t.font.name = '微軟正黑體'

    # Thin navy underline under title
    underline = slide.shapes.add_shape(
        1, inches(LP + 0.15), inches(0.73),
        inches(SW - LP - 0.3), inches(0.03))
    underline.fill.solid()
    underline.fill.fore_color.rgb = NAVY
    underline.line.fill.background()

    # Quote
    add_textbox(
        slide, LP + 0.15, 0.78, SW - LP - 0.3, 0.8,
        QUOTE,
        font_size=9, italic=True, color=NAVY,
        align=PP_ALIGN.LEFT, font_name='Times New Roman', wrap=True
    )

    # ── Puzzle tree image ─────────────────────────
    # Image: 4.2" × 5.88" placed at right side of right panel
    # figsize ratio 4.0:5.6 = 10:14, maintain proportions
    IMG_W = 4.2
    IMG_H = IMG_W * (5.6 / 4.0)   # ≈ 5.88" (maintain aspect 10:14)
    IMG_X = SW - IMG_W - 0.12     # ≈ 9.01"
    IMG_Y = 0.95

    tree_buf = draw_tree(
        show_roots=data['show_roots'],
        show_trunk=data['show_trunk'],
        show_canopy=data['show_canopy'],
    )
    slide.shapes.add_picture(tree_buf, inches(IMG_X), inches(IMG_Y),
                             inches(IMG_W), inches(IMG_H))

    # ── Stage label area (left side of right panel) ─────────────────
    stage = data['stage']
    accent_col = ACCENT[stage]
    TEXT_X = LP + 0.2
    TEXT_W = IMG_X - LP - 0.35   # width before the image

    # Stage circle indicator
    circle_size = 0.42
    circle_x = TEXT_X + 0.05
    circle_y = 1.35
    circ = slide.shapes.add_shape(
        9,   # MSO_SHAPE.OVAL
        inches(circle_x), inches(circle_y),
        inches(circle_size), inches(circle_size)
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = accent_col
    circ.line.fill.background()

    # Checkmark / sprout inside circle
    circ_txt = slide.shapes.add_textbox(
        inches(circle_x), inches(circle_y),
        inches(circle_size), inches(circle_size)
    )
    tf_c = circ_txt.text_frame
    p_c = tf_c.paragraphs[0]
    p_c.alignment = PP_ALIGN.CENTER
    run_c = p_c.add_run()
    run_c.text = '✓'
    run_c.font.size = Pt(14)
    run_c.font.bold = True
    run_c.font.color.rgb = WHITE_RGB
    run_c.font.name = 'Calibri'

    # Stage name (large)
    stage_box = slide.shapes.add_textbox(
        inches(TEXT_X + circle_size + 0.18), inches(1.28),
        inches(TEXT_W - circle_size - 0.25), inches(0.65)
    )
    tf_s = stage_box.text_frame
    p_s = tf_s.paragraphs[0]
    run_s = p_s.add_run()
    run_s.text = stage
    run_s.font.size = Pt(26)
    run_s.font.bold = True
    run_s.font.color.rgb = accent_col
    run_s.font.name = '微軟正黑體'

    # Sub-label
    add_textbox(
        slide,
        TEXT_X + circle_size + 0.18, 1.92,
        TEXT_W - circle_size - 0.25, 0.48,
        data['sub'],
        font_size=14, bold=True, color=NAVY,
        align=PP_ALIGN.LEFT, font_name='微軟正黑體'
    )

    # Horizontal divider below stage heading
    div = slide.shapes.add_shape(
        1, inches(TEXT_X), inches(2.46),
        inches(TEXT_W), inches(0.04)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = accent_col
    div.line.fill.background()

    # Bullet points
    bullet_x = TEXT_X + 0.05
    bullet_w = TEXT_W - 0.1
    bullets_box = slide.shapes.add_textbox(
        inches(bullet_x), inches(2.58),
        inches(bullet_w), inches(3.6)
    )
    tf_bul = bullets_box.text_frame
    tf_bul.word_wrap = True
    first = True
    for b in data['bullets']:
        if first:
            p_bul = tf_bul.paragraphs[0]
            first = False
        else:
            p_bul = tf_bul.add_paragraph()
        p_bul.space_before = Pt(6)
        run_bul = p_bul.add_run()
        run_bul.text = '●  ' + b
        run_bul.font.size = Pt(13)
        run_bul.font.color.rgb = DARK_NAVY
        run_bul.font.name = '微軟正黑體'

    # Bottom italic statement — spans text area only (left of tree image)
    bottom_y = FY - 0.58
    text_area_w = IMG_X - LP - 0.25
    add_textbox(
        slide, LP + 0.15, bottom_y, text_area_w, 0.52,
        data['bottom'],
        font_size=13, italic=True, bold=True,
        color=NAVY, align=PP_ALIGN.CENTER,
        font_name='微軟正黑體'
    )


# ─────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for data in SLIDE_DATA:
        build_slide(prs, data)

    out_path = '/workspace/puzzle_tree_presentation.pptx'
    prs.save(out_path)
    print(f'Saved: {out_path}')


if __name__ == '__main__':
    main()
